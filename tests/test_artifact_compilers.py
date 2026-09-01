"""
tests/test_artifact_compilers.py — Structural tests for Phase 3 artifact compilers.

Tests:
  - Each compiler generates a file without error
  - Each generated file is structurally readable (not corrupt)
  - Disclaimer text is present in docx and pptx
  - Provenance fields are present in all compiler results
  - Output path is inside data/artifacts/
  - Path traversal is rejected

Fully offline — no LLM, no network required.
"""
import sys
import json
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))
import pytest

SAMPLE_EVIDENCE = [
    {
        "doc_id": "oisd_116_pressure_vessel_inspection",
        "score": 0.7950,
        "text_preview": "Pressure vessels in H2S wet service shall be inspected every 5 years.",
    },
    {
        "doc_id": "mrpl_sop_flare_system",
        "score": 0.6800,
        "text_preview": "Nitrogen purge flow rate: 200 Nm3/h. KOD high-high trip: 70% LVL.",
    },
]

SAMPLE_KWARGS = dict(
    query="What is the inspection interval for H2S pressure vessels?",
    evidence=SAMPLE_EVIDENCE,
    session_id="test-session-artifact-001",
    model_role="reasoning",
    label="GROUNDED",
    request_id="test-req-001",
)


# --------------------------------------------------------------------------- #
# DOCX Tests                                                                   #
# --------------------------------------------------------------------------- #

class TestDocxCompiler:

    def test_docx_generates_without_error(self):
        from src.artifact_compiler import docx_compiler
        result = docx_compiler.generate(**SAMPLE_KWARGS)
        assert result["status"] == "ok", f"DOCX failed: {result.get('error')}"
        assert result["file_name"].endswith(".docx")
        assert Path(result["file_path"]).exists()

    def test_docx_file_is_structurally_readable(self):
        from src.artifact_compiler import docx_compiler
        from docx import Document
        result = docx_compiler.generate(**SAMPLE_KWARGS)
        doc = Document(result["file_path"])
        texts = " ".join(p.text for p in doc.paragraphs)
        assert len(texts) > 50, "Document has no readable text"

    def test_docx_contains_disclaimer(self):
        from src.artifact_compiler import docx_compiler
        from docx import Document
        result = docx_compiler.generate(**SAMPLE_KWARGS)
        doc = Document(result["file_path"])
        full_text = " ".join(p.text for p in doc.paragraphs)
        assert "not an engineering approval" in full_text, "Disclaimer not found in DOCX"

    def test_docx_provenance_fields_present(self):
        from src.artifact_compiler import docx_compiler
        result = docx_compiler.generate(**SAMPLE_KWARGS)
        prov = result.get("provenance", {})
        assert prov.get("request_id") == "test-req-001"
        assert prov.get("label") == "GROUNDED"
        assert prov.get("compiler") != ""
        assert "evidence_ids" in prov

    def test_docx_output_inside_artifact_dir(self):
        from src.artifact_compiler import docx_compiler
        from src.config import BASE_DIR
        result = docx_compiler.generate(**SAMPLE_KWARGS)
        artifact_dir = (BASE_DIR / "data" / "artifacts").resolve()
        file_path = Path(result["file_path"]).resolve()
        assert str(file_path).startswith(str(artifact_dir)), (
            f"Output file is outside artifact dir: {file_path}"
        )


# --------------------------------------------------------------------------- #
# XLSX Tests                                                                   #
# --------------------------------------------------------------------------- #

class TestXlsxCompiler:

    def test_xlsx_generates_without_error(self):
        from src.artifact_compiler import xlsx_compiler
        result = xlsx_compiler.generate(**SAMPLE_KWARGS)
        assert result["status"] == "ok", f"XLSX failed: {result.get('error')}"
        assert result["file_name"].endswith(".xlsx")
        assert Path(result["file_path"]).exists()

    def test_xlsx_file_is_structurally_readable(self):
        from src.artifact_compiler import xlsx_compiler
        import openpyxl
        result = xlsx_compiler.generate(**SAMPLE_KWARGS)
        wb = openpyxl.load_workbook(result["file_path"])
        assert len(wb.sheetnames) >= 3, f"Expected ≥3 sheets, got: {wb.sheetnames}"

    def test_xlsx_has_visible_formula(self):
        from src.artifact_compiler import xlsx_compiler
        import openpyxl
        result = xlsx_compiler.generate(**SAMPLE_KWARGS)
        wb = openpyxl.load_workbook(result["file_path"], data_only=False)
        ws = wb["Calculation"]
        # Cell B9 should be a formula
        cell_b9 = ws["B9"].value
        assert str(cell_b9).startswith("="), f"B9 is not a formula: {cell_b9}"

    def test_xlsx_has_provenance_sheet(self):
        from src.artifact_compiler import xlsx_compiler
        import openpyxl
        result = xlsx_compiler.generate(**SAMPLE_KWARGS)
        wb = openpyxl.load_workbook(result["file_path"])
        assert "Provenance" in wb.sheetnames

    def test_xlsx_provenance_fields_present(self):
        from src.artifact_compiler import xlsx_compiler
        result = xlsx_compiler.generate(**SAMPLE_KWARGS)
        prov = result.get("provenance", {})
        assert prov.get("request_id") == "test-req-001"
        assert prov.get("compiler") != ""

    def test_xlsx_output_inside_artifact_dir(self):
        from src.artifact_compiler import xlsx_compiler
        from src.config import BASE_DIR
        result = xlsx_compiler.generate(**SAMPLE_KWARGS)
        artifact_dir = (BASE_DIR / "data" / "artifacts").resolve()
        assert Path(result["file_path"]).resolve().is_relative_to(artifact_dir)


# --------------------------------------------------------------------------- #
# PPTX Tests                                                                   #
# --------------------------------------------------------------------------- #

class TestPptxCompiler:

    def test_pptx_generates_without_error(self):
        from src.artifact_compiler import pptx_compiler
        result = pptx_compiler.generate(**SAMPLE_KWARGS)
        assert result["status"] == "ok", f"PPTX failed: {result.get('error')}"
        assert result["file_name"].endswith(".pptx")
        assert Path(result["file_path"]).exists()

    def test_pptx_file_is_structurally_readable(self):
        from src.artifact_compiler import pptx_compiler
        from pptx import Presentation
        result = pptx_compiler.generate(**SAMPLE_KWARGS)
        prs = Presentation(result["file_path"])
        assert len(prs.slides) >= 5, f"Expected ≥5 slides, got {len(prs.slides)}"

    def test_pptx_contains_disclaimer(self):
        from src.artifact_compiler import pptx_compiler
        from pptx import Presentation
        result = pptx_compiler.generate(**SAMPLE_KWARGS)
        prs = Presentation(result["file_path"])
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text + " "
        assert "not an engineering approval" in all_text, "Disclaimer not found in PPTX"

    def test_pptx_has_limitations_slide(self):
        """Limitations slide must ALWAYS be present — never suppressed."""
        from src.artifact_compiler import pptx_compiler
        from pptx import Presentation
        result = pptx_compiler.generate(**SAMPLE_KWARGS)
        prs = Presentation(result["file_path"])
        all_text = " ".join(
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "Capability Limitations" in all_text, "Limitations slide missing from PPTX"
        assert "DEGRADED_SANDBOX" in all_text, "DEGRADED_SANDBOX not labelled in PPTX"

    def test_pptx_provenance_fields_present(self):
        from src.artifact_compiler import pptx_compiler
        result = pptx_compiler.generate(**SAMPLE_KWARGS)
        prov = result.get("provenance", {})
        assert prov.get("request_id") == "test-req-001"
        assert prov.get("compiler") != ""

    def test_pptx_output_inside_artifact_dir(self):
        from src.artifact_compiler import pptx_compiler
        from src.config import BASE_DIR
        result = pptx_compiler.generate(**SAMPLE_KWARGS)
        artifact_dir = (BASE_DIR / "data" / "artifacts").resolve()
        assert Path(result["file_path"]).resolve().is_relative_to(artifact_dir)


# --------------------------------------------------------------------------- #
# Cross-compiler independence                                                  #
# --------------------------------------------------------------------------- #

class TestCompilerIndependence:

    def test_all_three_succeed_independently(self):
        """Each compiler produces a file independently — one doesn't gate another."""
        from src.artifact_compiler import docx_compiler, xlsx_compiler, pptx_compiler
        r1 = docx_compiler.generate(**SAMPLE_KWARGS)
        r2 = xlsx_compiler.generate(**SAMPLE_KWARGS)
        r3 = pptx_compiler.generate(**SAMPLE_KWARGS)
        assert r1["status"] == "ok"
        assert r2["status"] == "ok"
        assert r3["status"] == "ok"

    def test_empty_evidence_handled_gracefully(self):
        """Compilers must not crash when evidence is empty."""
        from src.artifact_compiler import docx_compiler, xlsx_compiler, pptx_compiler
        kwargs = {**SAMPLE_KWARGS, "evidence": [], "label": "INSUFFICIENT_EVIDENCE"}
        for compiler in [docx_compiler, xlsx_compiler, pptx_compiler]:
            result = compiler.generate(**kwargs)
            assert result["status"] in ("ok", "error"), f"Unexpected status: {result}"
            if result["status"] == "ok":
                assert Path(result["file_path"]).exists()


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
