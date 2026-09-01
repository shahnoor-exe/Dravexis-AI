"""
artifact_compiler/pptx_compiler.py — Executive briefing PowerPoint generator.

Generates:
  - Title slide
  - Findings slide (from evidence)
  - Synthetic degradation curve slide (text-based, no matplotlib dependency)
  - Capability Limitations slide (always present — VISION_UNAVAILABLE, DEGRADED_SANDBOX)
  - Provenance slide
  - Disclaimer on every slide in footer

All outputs go to data/artifacts/ only.
"""
from __future__ import annotations

import logging
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from ..config import BASE_DIR

logger = logging.getLogger(__name__)

DISCLAIMER = (
    "Prototype output — not an engineering approval or statutory determination. "
    "MRPL Sovereign AI Workbench — Demo System."
)
ARTIFACT_DIR = BASE_DIR / "data" / "artifacts"
COMPILER_VERSION = "pptx_compiler-0.1.0 / python-pptx-1.0.2"

_MRPL_BLUE = RGBColor(0x1F, 0x38, 0x64)
_WARN_RED = RGBColor(0xC0, 0x00, 0x00)
_WARN_YELLOW = RGBColor(0xFF, 0xF2, 0xCC)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)


def _add_disclaimer_footer(slide, prs) -> None:
    """Add disclaimer text box at slide bottom."""
    left = Inches(0.3)
    top = prs.slide_height - Inches(0.45)
    width = prs.slide_width - Inches(0.6)
    height = Inches(0.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = DISCLAIMER
    run.font.size = Pt(7)
    run.font.color.rgb = _WARN_RED
    run.font.italic = True


def _add_slide(prs, layout_idx=6) -> Any:
    """Add a blank slide."""
    layout = prs.slide_layouts[layout_idx] if layout_idx < len(prs.slide_layouts) else prs.slide_layouts[-1]
    return prs.slides.add_slide(layout)


def _title_box(slide, prs, text: str, y=Inches(0.5), size=Pt(28), color=None) -> None:
    color = color or _MRPL_BLUE
    txBox = slide.shapes.add_textbox(Inches(0.5), y, prs.slide_width - Inches(1), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = color


def _body_box(slide, prs, text: str, y=Inches(1.7), height=Inches(4.5), size=Pt(14)) -> None:
    txBox = slide.shapes.add_textbox(Inches(0.5), y, prs.slide_width - Inches(1), height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = size


def generate(
    *,
    query: str,
    evidence: list[dict],
    session_id: str,
    model_role: str | None,
    label: str = "SYNTHETIC",
    request_id: str | None = None,
    vision_status: str = "VISION_UNAVAILABLE",
    sandbox_mode: str = "DEGRADED_SANDBOX",
) -> dict[str, Any]:
    ARTIFACT_DIR.mkdir(parents=True, exist_ok=True)
    request_id = request_id or str(uuid.uuid4())
    ts = datetime.now(timezone.utc)
    safe_name = f"briefing_{request_id[:8]}_{ts.strftime('%Y%m%d_%H%M%S')}.pptx"
    out_path = ARTIFACT_DIR / safe_name

    try:
        out_path.resolve().relative_to(ARTIFACT_DIR.resolve())
    except ValueError:
        return {"status": "error", "error": "path traversal rejected"}

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # ------------------------------------------------------------------ #
        # Slide 1: Title                                                       #
        # ------------------------------------------------------------------ #
        s1 = _add_slide(prs)
        # Background rectangle
        bg = s1.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(2.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = _MRPL_BLUE
        bg.line.fill.background()

        _title_box(s1, prs, "MRPL Sovereign AI Workbench", y=Inches(0.4), size=Pt(32), color=_WHITE)
        _title_box(s1, prs, "Equipment Inspection — Executive Briefing (DEMO)", y=Inches(1.2), size=Pt(20), color=_WHITE)
        _body_box(s1, prs,
            f"Query: {query[:120]}\n"
            f"Session: {session_id}\n"
            f"Generated: {ts.strftime('%Y-%m-%d %H:%M UTC')}\n"
            f"Label: {label}",
            y=Inches(2.8), size=Pt(13))
        _add_disclaimer_footer(s1, prs)

        # ------------------------------------------------------------------ #
        # Slide 2: Key Findings                                                #
        # ------------------------------------------------------------------ #
        s2 = _add_slide(prs)
        _title_box(s2, prs, "Key Findings", size=Pt(28))
        findings_text = ""
        if evidence:
            for ev in evidence[:4]:
                findings_text += (
                    f"• [{ev.get('doc_id', '?')}] Score: {ev.get('score', 0):.4f}\n"
                    f"  {ev.get('text_preview', '')[:200]}\n\n"
                )
        else:
            findings_text = "• No grounded evidence retrieved (INSUFFICIENT_EVIDENCE).\n• All findings below are synthetic demo data."
        _body_box(s2, prs, findings_text, size=Pt(12))
        _add_disclaimer_footer(s2, prs)

        # ------------------------------------------------------------------ #
        # Slide 3: Synthetic Degradation Curve                                 #
        # ------------------------------------------------------------------ #
        s3 = _add_slide(prs)
        _title_box(s3, prs, "Synthetic Degradation Curve (Demo Data)", size=Pt(24))
        curve_text = (
            "⚠  SYNTHETIC DATA — not from real inspection records.\n\n"
            "Year 0  → Thickness: 8.5 mm\n"
            "Year 1  → Thickness: 8.2 mm  (−0.3 mm)\n"
            "Year 2  → Thickness: 7.9 mm  (−0.3 mm)\n"
            "Year 3  → Thickness: 7.6 mm  (−0.3 mm)\n"
            "Year 4  → Thickness: 7.3 mm  (−0.3 mm)\n"
            "Year 5  → Thickness: 7.0 mm  (−0.3 mm)\n"
            "Year 6  → Thickness: 6.7 mm  (−0.3 mm)\n"
            "Year 7  → Thickness: 6.4 mm  (−0.3 mm)\n"
            "Year 8  → Thickness: 6.1 mm  (−0.3 mm)\n"
            "Year 8.3 → Threshold: 6.0 mm (ASME B31.3 minimum) ← Remaining Life Boundary\n\n"
            "Corrosion Rate (synthetic): 0.30 mm/year\n"
            "Remaining Life (synthetic): (8.5 − 6.0) / 0.30 = 8.33 years"
        )
        _body_box(s3, prs, curve_text, size=Pt(11))
        _add_disclaimer_footer(s3, prs)

        # ------------------------------------------------------------------ #
        # Slide 4: Capability Limitations (ALWAYS PRESENT)                    #
        # ------------------------------------------------------------------ #
        s4 = _add_slide(prs)
        # Warning background strip
        warn_bg = s4.shapes.add_shape(1, 0, Inches(0.3), prs.slide_width, Inches(0.9))
        warn_bg.fill.solid()
        warn_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xCC)
        warn_bg.line.fill.background()

        _title_box(s4, prs, "⚠  Capability Limitations — Must Read", size=Pt(24), color=_WARN_RED)
        limitations = (
            f"Vision Status: {vision_status}\n"
            "  → VL model (Qwen2.5-VL-3B) not loaded. No visual analysis performed.\n"
            "  → P&ID/image descriptions are not available in this demo run.\n\n"
            f"Sandbox Mode: {sandbox_mode}\n"
            "  → Docker not installed. Code runs in-process with allowlist only.\n"
            "  → NOT equivalent to Docker container isolation. Labelled DEGRADED.\n\n"
            "Corpus: 9 vectors (minimal demo set)\n"
            "  → Not production coverage. Expand before field deployment.\n\n"
            "Model Swap Latency: UNMEASURED\n"
            "  → GGUF model files not yet downloaded. Latency budgeting incomplete.\n\n"
            "This slide is always included in generated briefings. Do not suppress."
        )
        _body_box(s4, prs, limitations, y=Inches(1.4), size=Pt(11))
        _add_disclaimer_footer(s4, prs)

        # ------------------------------------------------------------------ #
        # Slide 5: Provenance                                                  #
        # ------------------------------------------------------------------ #
        s5 = _add_slide(prs)
        _title_box(s5, prs, "Provenance & Traceability", size=Pt(24))
        prov_text = (
            f"Request ID:     {request_id}\n"
            f"Session ID:     {session_id}\n"
            f"Timestamp UTC:  {ts.isoformat()}\n"
            f"Data Label:     {label}\n"
            f"Model Role:     {model_role or 'N/A'}\n"
            f"Compiler:       {COMPILER_VERSION}\n"
            f"Evidence Count: {len(evidence)}\n"
            f"Evidence IDs:   {', '.join(ev.get('doc_id','?') for ev in evidence[:5]) or 'none'}\n\n"
            f"Disclaimer: {DISCLAIMER}"
        )
        _body_box(s5, prs, prov_text, size=Pt(11))
        _add_disclaimer_footer(s5, prs)

        prs.save(str(out_path))
        logger.info("PPTX generated: %s", out_path.name)

    except Exception as exc:
        logger.error("PPTX generation failed: %s", exc)
        return {"status": "error", "error": str(exc)}

    provenance = {
        "request_id": request_id,
        "session_id": session_id,
        "timestamp_utc": ts.isoformat(),
        "evidence_ids": [ev.get("doc_id") for ev in evidence],
        "label": label,
        "model_role": model_role,
        "compiler": COMPILER_VERSION,
        "output_file": safe_name,
    }
    return {
        "status": "ok",
        "file_path": str(out_path),
        "file_name": safe_name,
        "provenance": provenance,
    }
